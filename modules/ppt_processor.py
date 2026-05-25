"""PPT/PPTX portfolio slide content extraction."""

from __future__ import annotations

from pathlib import Path


def extract_ppt_text(file_path: str | Path) -> dict:
    """Extract text from presentation slides."""
    from pptx import Presentation

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPT not found: {path}")

    try:
        prs = Presentation(str(path))
        slides_content = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides_content.append(f"Slide {idx}:\n" + "\n".join(slide_text))
        full_text = "\n\n".join(slides_content)
    except Exception as exc:
        raise RuntimeError(f"PPT extraction failed: {exc}") from exc

    return {
        "text": full_text,
        "metadata": {
            "source": path.name,
            "type": "portfolio_ppt",
            "slides": len(prs.slides),
            "char_count": len(full_text),
        },
    }
